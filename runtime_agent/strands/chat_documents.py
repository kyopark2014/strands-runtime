"""Document loading and summarization helpers for the chat module."""

import base64
import csv
import logging
import os
import traceback
import uuid
from io import BytesIO
from urllib import parse

import PyPDF2
from PIL import Image
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_core.prompts import ChatPromptTemplate
from langchain_text_splitters import RecursiveCharacterTextSplitter

from chat_s3 import get_s3_client, get_s3_resource

logger = logging.getLogger("chat")

SUMMARY_MAX_CHARS = 500
MIN_EXTRACTED_TEXT_LENGTH = 10
MAX_LLM_RETRY_ATTEMPTS = 5

# Cap decoded image dimensions before Bedrock multimodal encode: keeps memory and
# base64 payload size bounded (Bedrock image input practical limit ~2M pixels).
MAX_IMAGE_PIXELS = 2_000_000

IMAGE_FILE_TYPES = ("png", "jpeg", "jpg", "webp", "gif")
DOCUMENT_FILE_TYPES = ("pdf", "txt", "md", "pptx", "docx")


class DocumentSummaryService:
    """LLM-based text summarization for document content."""

    def build_prompt(self, text: str) -> ChatPromptTemplate:
        import chat

        if chat.isKorean(text) == True:
            system = (
                f"다음의 <article> tag안의 문장을 요약해서 {SUMMARY_MAX_CHARS}자 이내로 설명하세오."
            )
        else:
            system = (
                "Here is pieces of article, contained in <article> tags. "
                f"Write a concise summary within {SUMMARY_MAX_CHARS} characters."
            )

        human = "<article>{text}</article>"
        return ChatPromptTemplate.from_messages([("system", system), ("human", human)])

    def invoke(self, docs) -> str:
        import chat

        llm = chat.get_chat(extended_thinking=chat.reasoning_mode)

        text = ""
        for doc in docs:
            text = text + doc

        prompt = self.build_prompt(text)
        chain = prompt | llm
        try:
            result = chain.invoke({"text": text})

            summary = result.content
            logger.info(f"esult of summarization: {summary}")
        except Exception:
            err_msg = traceback.format_exc()
            logger.info(f"error message: {err_msg}")
            raise Exception("Not able to request to LLM")

        return summary


class DocumentLoaderService:
    """S3 document retrieval, format parsing, and text splitting."""

    def fetch_s3_object(self, s3_file_name):
        import chat

        s3r = get_s3_resource()
        doc = s3r.Object(chat.s3_bucket, chat.s3_prefix + "/" + s3_file_name)
        logger.info(
            f"s3_bucket: {chat.s3_bucket}, s3_prefix: {chat.s3_prefix}, "
            f"s3_file_name: {s3_file_name}"
        )
        return doc

    def parse_pdf(self, raw_bytes: bytes) -> str:
        reader = PyPDF2.PdfReader(BytesIO(raw_bytes))

        raw_text = []
        for page in reader.pages:
            raw_text.append(page.extract_text())
        return "\n".join(raw_text)

    def parse_text_file(self, raw_bytes: bytes) -> str:
        return raw_bytes.decode("utf-8")

    def split_text(self, contents: str) -> list:
        logger.info(f"contents: {contents}")
        new_contents = str(contents).replace("\n", " ")
        logger.info(f"length: {len(new_contents)}")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function=len,
        )
        texts = text_splitter.split_text(new_contents)
        if texts:
            logger.info(f"exts[0]: {texts[0]}")

        return texts

    def load(self, file_type, s3_file_name):
        doc = self.fetch_s3_object(s3_file_name)

        contents = ""
        if file_type == "pdf":
            contents = self.parse_pdf(doc.get()["Body"].read())
        elif file_type == "txt" or file_type == "md":
            contents = self.parse_text_file(doc.get()["Body"].read())

        return self.split_text(contents)

    def load_csv(self, s3_file_name):
        doc = self.fetch_s3_object(s3_file_name)

        lines = doc.get()["Body"].read().decode("utf-8").split("\n")  # read csv per line
        logger.info(f"lins: {len(lines)}")

        columns = lines[0].split(",")  # get columns
        logger.info(f"columns: {columns}")

        docs = []
        n = 0
        for row in csv.DictReader(lines, delimiter=",", quotechar='"'):
            values = {k: row[k] for k in columns if k in row}
            content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
            doc = Document(
                page_content=content,
                metadata={
                    "name": s3_file_name,
                    "row": n + 1,
                }
            )
            docs.append(doc)
            n = n + 1
        logger.info(f"docs[0]: {docs[0]}")

        return docs


_summary_service = DocumentSummaryService()
_loader_service = DocumentLoaderService()


def get_summary(docs):
    return _summary_service.invoke(docs)


def load_document(file_type, s3_file_name):
    return _loader_service.load(file_type, s3_file_name)


def _file_name_from_ref(file_ref: str) -> str:
    """Extract a basename from a sharing URL or plain file name."""
    raw = (file_ref or "").strip()
    if not raw:
        return ""
    name = raw.rsplit("/", 1)[-1]
    return parse.unquote(name)


def _s3_key_from_file_ref(file_ref: str, *, default_prefix: str) -> str:
    """Derive S3 object key from a sharing URL or bare file name.

    Prefer the ``images/{user_id}/...`` (or docs/) path embedded in the URL;
    otherwise fall back to ``{prefix}/{user_id}/{file_name}``.
    """
    import chat
    from tools.workspace import sanitize_user_path_segment

    raw = (file_ref or "").strip()
    if not raw:
        return ""

    raw = raw.split("?", 1)[0].split("#", 1)[0]

    for prefix in (chat.s3_image_prefix, chat.s3_prefix, "images", "docs"):
        marker = f"/{prefix}/"
        idx = raw.find(marker)
        if idx >= 0:
            return parse.unquote(raw[idx + 1 :])
        if raw.startswith(f"{prefix}/"):
            return parse.unquote(raw)

    file_name = _file_name_from_ref(file_ref)
    if not file_name:
        return ""

    user_segment = sanitize_user_path_segment(chat.user_id)
    if user_segment:
        return f"{default_prefix}/{user_segment}/{file_name}"
    return f"{default_prefix}/{file_name}"


def _workspace_ref_to_s3_key(file_ref: str) -> str | None:
    """Map /mnt/workspace/{user}/upload/x → agentcore-sessions/{user}/upload/x."""
    path = (file_ref or "").strip()
    marker = "/mnt/workspace/"
    if not path.startswith(marker):
        return None
    rel = path[len(marker) :].lstrip("/")
    if not rel or ".." in rel.split("/"):
        return None
    return f"agentcore-sessions/{rel}"


def _wait_for_workspace_mount_file(
    path: str,
    *,
    timeout_sec: float = 90.0,
    interval_sec: float = 1.0,
) -> bool:
    """Poll until ``path`` appears under /mnt/workspace (S3 Files lag)."""
    import time

    if not path.startswith("/mnt/workspace/"):
        return False
    deadline = time.monotonic() + max(0.0, timeout_sec)
    while True:
        if os.path.isfile(path) and os.path.getsize(path) > 0:
            logger.info(f"workspace mount file ready: {path}")
            return True
        if time.monotonic() >= deadline:
            logger.warning(
                f"workspace mount file not visible after {timeout_sec:.0f}s: {path}"
            )
            return False
        time.sleep(max(0.1, interval_sec))


def _load_workspace_file_bytes(file_ref: str) -> tuple[bytes | None, str]:
    """Load Load-files bytes from /mnt/workspace, waiting for mount if needed.

    Falls back to S3 API only if the mount never shows the file.
    """
    import chat

    path = (file_ref or "").strip()
    if path.startswith("/mnt/workspace/") and os.path.isfile(path):
        with open(path, "rb") as f:
            data = f.read()
        logger.info(f"loaded workspace file from mount ({len(data)} bytes): {path}")
        return data, path

    if path.startswith("/mnt/workspace/"):
        logger.info(f"waiting for workspace mount file: {path}")
        if _wait_for_workspace_mount_file(path, timeout_sec=90.0, interval_sec=1.0):
            with open(path, "rb") as f:
                data = f.read()
            logger.info(
                f"loaded workspace file after wait ({len(data)} bytes): {path}"
            )
            return data, path

    s3_key = _workspace_ref_to_s3_key(path)
    if not s3_key or not chat.s3_bucket:
        logger.warning(
            "workspace file unavailable on mount and no S3 key/bucket: path=%s key=%s",
            path,
            s3_key,
        )
        return None, path

    try:
        s3_client = get_s3_client()
        logger.info(f"loading workspace file from s3://{chat.s3_bucket}/{s3_key}")
        obj = s3_client.get_object(Bucket=chat.s3_bucket, Key=s3_key)
        data = obj["Body"].read()
        logger.info(f"loaded workspace file from S3 ({len(data)} bytes): {s3_key}")
        try:
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "wb") as f:
                f.write(data)
        except Exception as cache_err:
            logger.warning(f"could not cache S3 file onto mount {path}: {cache_err}")
        return data, f"s3://{chat.s3_bucket}/{s3_key}"
    except Exception:
        logger.error(
            "Failed to load workspace file from S3 key=%s: %s",
            s3_key,
            traceback.format_exc(),
        )
        return None, path


def _extract_text_from_docx_bytes(data: bytes) -> str:
    try:
        import docx

        document = docx.Document(BytesIO(data))
        parts: list[str] = []
        for p in document.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text)
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            return text
    except Exception as e:
        logger.warning(f"python-docx extract failed, trying zip/xml: {e}")

    import zipfile
    from xml.etree import ElementTree as ET

    with zipfile.ZipFile(BytesIO(data)) as zf:
        xml_bytes = zf.read("word/document.xml")
    root = ET.fromstring(xml_bytes)
    paragraphs: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}p") or node.tag == "p":
            runs = [
                (t.text or "")
                for t in node.iter()
                if (t.tag.endswith("}t") or t.tag == "t") and t.text
            ]
            line = "".join(runs).strip()
            if line:
                paragraphs.append(line)
    if paragraphs:
        return "\n".join(paragraphs)
    texts = [
        (node.text or "")
        for node in root.iter()
        if node.tag.endswith("}t") or node.tag == "t"
    ]
    return "\n".join(t for t in texts if t).strip()


def _extract_text_from_legacy_doc_bytes(data: bytes) -> str:
    """Best-effort extraction for Word 97-2003 (.doc)."""
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        doc_path = os.path.join(tmp, "input.doc")
        with open(doc_path, "wb") as f:
            f.write(data)

        for cmd in (
            ["antiword", doc_path],
            ["catdoc", "-w", doc_path],
        ):
            try:
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=60,
                    check=False,
                )
                text = (result.stdout or "").strip()
                if result.returncode == 0 and text:
                    return text
            except FileNotFoundError:
                continue
            except Exception as e:
                logger.warning("legacy .doc extractor %s failed: %s", cmd[0], e)

        # LibreOffice → txt
        try:
            result = subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    tmp,
                    doc_path,
                ],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            txt_path = os.path.join(tmp, "input.txt")
            if os.path.isfile(txt_path):
                with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                    text = f.read().strip()
                if text:
                    return text
            if result.returncode != 0:
                logger.warning(
                    "soffice .doc convert failed: %s",
                    (result.stderr or result.stdout or "")[:500],
                )
        except FileNotFoundError:
            pass
        except Exception as e:
            logger.warning("soffice .doc convert failed: %s", e)

    # Last resort: pull printable / UTF-16LE runs from the OLE binary.
    try:
        import re

        texts: list[str] = []
        for offset in (0, 1):
            for match in re.finditer(rb"(?:[\x20-\x7e]\x00){6,}", data[offset:]):
                try:
                    chunk = match.group().decode("utf-16le", errors="ignore").strip()
                except Exception:
                    continue
                if len(chunk) >= 6 and chunk not in texts:
                    texts.append(chunk)
        ascii_runs = re.findall(rb"[\x20-\x7e]{8,}", data)
        for raw in ascii_runs:
            chunk = raw.decode("ascii", errors="ignore").strip()
            if len(chunk) >= 8 and chunk not in texts:
                texts.append(chunk)
        joined = "\n".join(texts).strip()
        if len(joined) >= 40:
            return joined
    except Exception as e:
        logger.warning("legacy .doc binary string extract failed: %s", e)

    return ""


def _extract_text_from_pdf_bytes(data: bytes) -> str:
    try:
        import pdfplumber

        parts: list[str] = []
        with pdfplumber.open(BytesIO(data)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)
        text = "\n\n".join(parts).strip()
        if text:
            return text
    except Exception as e:
        logger.warning("pdfplumber bytes extract failed (%s); trying pypdf", e)

    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        parts = [(page.extract_text() or "") for page in reader.pages]
        return "\n\n".join(parts).strip()
    except Exception as e:
        raise RuntimeError(f"PDF 텍스트 추출 실패: {e}") from e


def _extract_text_from_document_bytes(data: bytes, file_type: str) -> str:
    """Best-effort text extraction for Load-files (bytes from mount or S3)."""
    if file_type in (
        "txt",
        "md",
        "markdown",
        "csv",
        "json",
        "py",
        "js",
        "ts",
        "tsx",
        "jsx",
        "html",
        "htm",
        "yml",
        "yaml",
        "xml",
        "rst",
    ):
        return data.decode("utf-8", errors="replace")

    if file_type == "pdf":
        return _extract_text_from_pdf_bytes(data)

    if file_type == "docx":
        return _extract_text_from_docx_bytes(data)

    if file_type == "doc":
        return _extract_text_from_legacy_doc_bytes(data)

    if file_type == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(BytesIO(data))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if text and str(text).strip():
                        parts.append(str(text).strip())
            return "\n".join(parts).strip()
        except Exception as e:
            raise RuntimeError(f"PPTX 텍스트 추출 실패: {e}") from e

    if file_type in ("xlsx", "xls"):
        try:
            import openpyxl

            wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(v.strip() for v in vals):
                        rows.append("\t".join(vals))
            return "\n".join(rows).strip()
        except Exception as e:
            raise RuntimeError(f"Excel 텍스트 추출 실패: {e}") from e

    return ""


def _extract_text_from_local_document(file_path: str, file_type: str) -> str:
    """Best-effort text extraction for Load-files under /mnt/workspace."""
    with open(file_path, "rb") as f:
        data = f.read()
    return _extract_text_from_document_bytes(data, file_type)


def summarize_image_file(file_ref, prompt=""):
    import chat

    s3_client = get_s3_client()
    s3_key = _s3_key_from_file_ref(file_ref, default_prefix=chat.s3_image_prefix)
    logger.info(f"loading image from s3://{chat.s3_bucket}/{s3_key}")
    try:
        image_obj = s3_client.get_object(Bucket=chat.s3_bucket, Key=s3_key)
        image_content = image_obj["Body"].read()
    except Exception:
        logger.exception("Failed to load image from S3 key=%s", s3_key)
        return "이미지 파일을 불러오지 못했습니다."

    image_summary_prompt = (
        "사용자의 요청을 참조하여 이미지의 내용을 분석한 후에 markdown 포맷으로 "
        "자세히 설명해주세요. 사용자 요청: <user_request>"
        f"{prompt}</user_request>"
    )
    logger.info(f"image_summary_prompt: {image_summary_prompt}")

    return summarize_image(image_content, image_summary_prompt)


def summarize_csv_file(file_name):
    docs = load_csv_document(file_name)
    contexts = []
    for doc in docs:
        contexts.append(doc.page_content)
    logger.info(f"contexts: {contexts}")

    return get_summary(contexts)


def summarize_document_file(file_name, file_type):
    import chat

    texts = load_document(file_type, file_name)

    if len(texts):
        docs = []
        for i in range(len(texts)):
            docs.append(
                Document(
                    page_content=texts[i],
                    metadata={
                        "name": file_name,
                        # 'page':i+1,
                        "url": chat.path
                        + "/"
                        + chat.doc_prefix
                        + parse.quote(file_name),
                    },
                )
            )
        logger.info(f"docs[0]: {docs[0]}")
        logger.info(f"docs size: {len(docs)}")

        contexts = []
        for doc in docs:
            contexts.append(doc.page_content)
        logger.info(f"contexts: {contexts}")

        return get_summary(contexts)

    return "문서 로딩에 실패하였습니다."


def get_summary_of_uploaded_file(file_ref, st=None, prompt=""):
    """Analyze an uploaded file (by URL, workspace path, or name) and return a text summary.

    Images are loaded from S3 under images/{user_id}/ and summarized with vision.
    Load-files paths under /mnt/workspace/{user}/upload/ are read from the
    session-storage mount, with an S3 API fallback when the mount lags.
    Documents keep the existing docs/ load path.
    """
    import chat

    file_name = _file_name_from_ref(file_ref)
    if not file_name:
        return "파일 이름을 확인할 수 없습니다."

    file_type = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    logger.info(
        f"get_summary_of_uploaded_file: file_name={file_name}, file_type={file_type}"
    )

    # Load-files: agentcore-sessions → /mnt/workspace/{user}/upload/{name}
    local_path = (file_ref or "").strip()
    if local_path.startswith("/mnt/workspace/"):
        data, source = _load_workspace_file_bytes(local_path)
        if data is None:
            return (
                f"파일 경로: {local_path}\n"
                "파일을 마운트/S3에서 읽지 못했습니다. "
                f"read_file 또는 bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )

        if file_type in IMAGE_FILE_TYPES:
            image_summary_prompt = (
                "사용자의 요청을 참조하여 이미지의 내용을 분석한 후에 markdown 포맷으로 자세히 설명해주세요. "
                f"사용자 요청: <user_request>{prompt}</user_request>"
            )
            return summarize_image(data, image_summary_prompt)

        try:
            extracted = _extract_text_from_document_bytes(data, file_type)
        except Exception as e:
            logger.error(f"Failed to extract text from {local_path} ({source}): {e}")
            return (
                f"파일 경로: {local_path}\n"
                f"텍스트 추출에 실패했습니다 ({e}). "
                f"read_file / bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )
        if not extracted:
            return (
                f"파일 경로: {local_path}\n"
                f"파일 형식(.{file_type})에서 텍스트를 추출하지 못했습니다. "
                f"read_file / bash로 `{local_path}` 를 직접 읽어 분석하세요."
            )
        max_chars = 100000
        if len(extracted) > max_chars:
            extracted = extracted[:max_chars] + "\n\n…(이하 생략)"
        return f"파일 경로: {local_path}\n\n{extracted}"

    if file_type in IMAGE_FILE_TYPES:
        return summarize_image_file(file_ref, prompt)

    msg = "지원하지 않는 파일 형식입니다."

    if file_type == "csv":
        msg = summarize_csv_file(file_name)
    elif file_type in DOCUMENT_FILE_TYPES:
        msg = summarize_document_file(file_name, file_type)

    chat.fileId = uuid.uuid4().hex

    return msg


def load_csv_document(s3_file_name):
    return _loader_service.load_csv(s3_file_name)


def _resize_and_encode(image_content):
    """Resize image if needed and return base64-encoded string."""
    img = Image.open(BytesIO(image_content))
    width, height = img.size
    logger.info(f"width: {width}, height: {height}, size: {width*height}")

    isResized = False
    max_size = 5 * 1024 * 1024  # 5MB

    while width * height > MAX_IMAGE_PIXELS:
        width = int(width / 2)
        height = int(height / 2)
        isResized = True

    if isResized:
        img = img.resize((width, height))

    max_attempts = 5
    img_base64 = ""
    base64_size = 0
    for attempt in range(max_attempts):
        buffer = BytesIO()
        img.save(buffer, format="PNG", optimize=True)
        img_bytes = buffer.getvalue()
        img_base64 = base64.b64encode(img_bytes).decode("utf-8")

        base64_size = len(img_base64.encode("utf-8"))
        logger.info(f"attempt {attempt + 1}: base64_size = {base64_size} bytes")

        if base64_size <= max_size:
            break
        width = int(width * 0.8)
        height = int(height * 0.8)
        img = img.resize((width, height))
        logger.info(f"resizing to {width}x{height} due to size limit")

    if base64_size > max_size:
        raise Exception("이미지 크기가 너무 큽니다. 5MB 이하의 이미지를 사용해주세요.")

    return img_base64


def extract_text(img_base64):
    """Extract text from an image using multimodal LLM."""
    import chat

    multimodal = chat.get_chat(extended_thinking=chat.reasoning_mode)
    query = "텍스트를 추출해서 markdown 포맷으로 변환하세요. <result> tag를 붙여주세요."

    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}",
                    },
                },
                {"type": "text", "text": query},
            ]
        )
    ]

    extracted_text = ""
    for attempt in range(MAX_LLM_RETRY_ATTEMPTS):
        logger.info(f"extract_text attempt: {attempt}")
        try:
            result = multimodal.invoke(messages)
            extracted_text = result.content
            break
        except Exception:
            err_msg = traceback.format_exc()
            logger.info(f"error message: {err_msg}")

    logger.info(f"extracted_text: {extracted_text}")
    if len(extracted_text) < MIN_EXTRACTED_TEXT_LENGTH:
        extracted_text = "텍스트를 추출하지 못하였습니다."

    return extracted_text


def summary_image(img_base64, instruction):
    """Summarize an image using multimodal LLM."""
    import chat

    llm = chat.get_chat(extended_thinking=chat.reasoning_mode)

    if instruction:
        logger.info(f"instruction: {instruction}")
        query = f"{instruction}. <result> tag를 붙여주세요. 한국어로 답변하세요."
    else:
        query = (
            "이미지가 의미하는 내용을 풀어서 자세히 알려주세요. "
            "markdown 포맷으로 답변을 작성합니다."
        )

    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}",
                    },
                },
                {"type": "text", "text": query},
            ]
        )
    ]

    for attempt in range(MAX_LLM_RETRY_ATTEMPTS):
        logger.info(f"summary_image attempt: {attempt}")
        try:
            result = llm.invoke(messages)
            extracted_text = result.content
            break
        except Exception:
            err_msg = traceback.format_exc()
            logger.info(f"error message: {err_msg}")
            raise Exception("Not able to request to LLM")

    return extracted_text


def summarize_image(image_content: bytes, prompt: str) -> str:
    """Resize image and summarize with vision."""
    img_base64 = _resize_and_encode(image_content)

    logger.info("이미지의 내용을 분석합니다.")
    result = summary_image(img_base64, prompt)

    summary = result[result.find("<result>") + 8 : result.find("</result>")]
    logger.info(f"image summary: {summary}")

    return summary
